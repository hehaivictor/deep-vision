import pptx
import sys

def read_presentation(file_path):
    prs = pptx.Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for j, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            print(f"Shape {j}: type={shape.shape_type}, name={shape.name}, text='{shape.text[:50]}'")

if __name__ == "__main__":
    read_presentation(sys.argv[1])
