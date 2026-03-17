import pptx
from pptx.util import Inches, Pt
import sys
import os

def create_textbox(slide, texts):
    left = Inches(1.0)
    top = Inches(1.8)
    width = Inches(11.33)
    height = Inches(5.2)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_info in enumerate(texts):
        text, level, bold = line_info
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = text
        p.level = level
        # p.font.name = 'Microsoft YaHei' 或者 'Arial'
        p.font.size = Pt(20 if level == 1 else 24)
        p.font.bold = bold

def modify_presentation(file_path):
    prs = pptx.Presentation(file_path)

    # 预定义每页的内容。结构: (文本, level, bold)
    slide2_texts = [
        ("一、 应用场景", 0, True),
        ("需求调研：面向客户方的 IT 项目、产品需求访谈收集（例如：企业级软件的需求定义）。", 1, False),
        ("业务分析：深度的业务模式拆解与竞品分析、技术方案评估。", 1, False),
        ("面试评估：针对于专业人才的系统化提问和面试考察。", 1, False),
        ("", 0, False),
        ("二、 核心需求（痛点）", 0, True),
        ("访谈无序与遗漏：人工提问缺乏系统框架，容易遗漏关键维度（如商业模式、技术挑战、性能指标等）。", 1, False),
        ("记录难与沉淀慢：访谈过程中记录信息不全，事后整理结构化文档极度耗时。", 1, False),
        ("追问深度不够：对专业领域的回答缺乏即时的深度追问能力，难以获取足够支撑最终报告的“实证级别”信息。", 1, False),
        ("交付转化低效：从访谈记录直接生成标准、高质量、有排版格式的交付报告流程繁琐冗长。", 1, False)
    ]

    slide3_texts = [
        ("一、 应用效果", 0, True),
        ("全链路 AI 引导：实现从智能提问、自动引导追问、用户回答实时沉淀到一键生成专业报告的一站式闭环。", 1, False),
        ("专业度与效率双升：内置场景和灵活自定义场景模型驱动，将资深咨询师/需求分析师的访谈思路产品化。", 1, False),
        ("高质量交付：异步生成包含摘要、主体评估、行动建议在内的完整结构化报告，极大减少人工案头工作。", 1, False),
        ("", 0, False),
        ("二、 核心界面与功能模块", 0, True),
        ("安全隔离的登录与场景选择：手机号验证码/微信扫码安全验证；提供内置场景库与自定义场景管理。", 1, False),
        ("场景化智能访谈台：呈现渐进式、多轮次、基于“探索与取证”双重策略的动态提问与对话界面。", 1, False),
        ("进度与状态反馈栏：实时展示当前访谈完整度与已覆盖访谈维度，引导用户推进对话进度。", 1, False),
        ("报告预览与管理台：清晰展示异步生成排队进度，支持多格式一键导出（含 Markdown、DOCX 及附录 PDF）。", 1, False)
    ]

    slide4_texts = [
        ("一、 技术栈体系", 0, True),
        ("后端架构：基于 Python + Flask（Gunicorn 生产部署），轻量化无状态服务架构，高可维护性。", 1, False),
        ("前端呈现：原生 HTML/CSS/JS 组合（内置 Alpine.js, Tailwind, Markdown 解析），极简部署高速响应。", 1, False),
        ("", 0, False),
        ("二、 AI 与核心逻辑设计", 0, True),
        ("多模型路由协同：细分任务派发最佳模型。提问（Minimax-2.5）、草稿（Kimi-k2.5）、审阅与摘要评分（GLM-5）。", 1, False),
        ("会话状态管理：基于 ETag 缓存控制、持久化本地结构化目录存储，及元数据索引回退机制进行状态校验与并发保护。", 1, False),
        ("异步队列引擎：基于独立工作队列与线程/进程池处理长耗时报告生成任务，支持最大排队控制（Max Pending）。", 1, False),
        ("多重并发防护：通过状态轮询、API 过载保护及 429（Too Many Requests）快速失败机制，保障系统高可用。", 1, False)
    ]

    slide5_texts = [
        ("一、 近期迭代与重构计划（V3 报告链路升级）", 0, True),
        ("第一阶段（先止血优化）：优化问题生成基础策略，去除导致误伤的生硬规则，解决错误的提问链路竞争（如 summary lane）。", 1, False),
        ("第二阶段（指标再校准）：重构访谈衡量口径，建立“覆盖率（关键方面+质量加权）”和“弱绑定比例（Weak Binding）”双指标评估。", 1, False),
        ("第三阶段（双轨制重构）：引入“探索型（探索方向收敛） vs 取证型（详实论据和场景）”双模式提问架构，高密度提取实证信息。", 1, False),
        ("第四阶段（交付与闭环）：全面上线 V3 报告生成系统，结合自动化回放测试机制验证质量收益与线上性能表现。", 1, False),
        ("", 0, False),
        ("二、 长期规划", 0, True),
        ("丰富内置与自定义的深度行业场景库（如金融风控、医疗信息化调研等），完善多租户数据沙箱隔离与多工作树协作发布机制。", 1, False)
    ]

    slide6_texts = [
        ("一、 技术挑战与应对", 0, True),
        ("大模型时延与稳定性：多模型串行/并行中单点超时导致全链路卡顿。需进一步优化请求超时策略、降级方案与快速重试机制。", 1, False),
        ("上下文长度与记忆衰减：复杂业务深研会累积超长对话历史，需研究更精细化的 Context 窗口滑动管理与记忆压缩摘要算法。", 1, False),
        ("高并发下的模型限流（Rate Limit）：多会话并发生成万字报告时，极易触碰 LLM API 速率限制，需增强全局熔断与分布式排队控制。", 1, False),
        ("", 0, False),
        ("二、 业务与产品关注点", 0, True),
        ("提问与报告的一致性：通过“探索+取证”模式升级后，依然需要关注前期提问导向对最终报告深度的真实转化率。", 1, False),
        ("数据隐私与合规：多租户/多实例混合部署下（INSTANCE_SCOPE_KEY），安全审计、账号权限归属权迁移与隔离合规机制需持续加固。", 1, False)
    ]

    # Slide 索引是 1, 2, 3, 4, 5 （对应第 2 到 第 6 页）
    if len(prs.slides) >= 6:
        create_textbox(prs.slides[1], slide2_texts)
        create_textbox(prs.slides[2], slide3_texts)
        create_textbox(prs.slides[3], slide4_texts)
        create_textbox(prs.slides[4], slide5_texts)
        create_textbox(prs.slides[5], slide6_texts)

    # 保存文件
    new_file_path = file_path.replace(".pptx", "_Filled.pptx")
    prs.save(new_file_path)
    print(f"Successfully saved filled presentation to: {new_file_path}")
    
    # 用覆盖的方式，将新文件移回原路径
    os.replace(new_file_path, file_path)
    print(f"Successfully replaced original presentation.")

if __name__ == "__main__":
    modify_presentation(sys.argv[1])
