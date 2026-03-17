import pptx
import sys

def read_presentation(file_path):
    try:
        prs = pptx.Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            print(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                print(f"[{shape.shape_type}] {shape.text}")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        read_presentation(sys.argv[1])
    else:
        print("Provide file path")
